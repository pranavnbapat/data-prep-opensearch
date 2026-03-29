# enricher_utils.py

import json
import ipaddress
import logging
import os
import random
import re
import socket
import subprocess
import threading
import time

from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import urlparse
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import requests

from pipeline.io import resolve_latest_pointer, find_latest_matching
from stages.downloader.utils import sha256_obj, stable_str
from requests.auth import HTTPBasicAuth


_tls = threading.local()
logger = logging.getLogger(__name__)
_RESOLVE_CACHE: Dict[str, str] = {}
_RESOLVE_LOCK = threading.Lock()
_OOXML_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
_SAFE_BROWSING_CACHE: Dict[str, bool] = {}
_SAFE_BROWSING_LOCK = threading.Lock()
_DANGEROUS_EXTENSIONS = {
    ".exe", ".msi", ".bat", ".cmd", ".ps1", ".vbs", ".dll",
    ".pkg", ".dmg", ".app",
    ".deb", ".rpm", ".appimage", ".sh", ".run",
    ".jar", ".py", ".pyz", ".pyc",
    ".scr", ".com", ".pif", ".msix", ".msixbundle", ".reg",
    ".iso", ".img", ".bin", ".apk", ".ipa",
    ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz", ".tgz", ".tbz2", ".txz",
}
_LEGACY_WEBPAGE_EXTENSIONS = {".php", ".cgi", ".jsp", ".asp", ".aspx", ".cfm"}
_SHORTENER_HOSTS = {
    "bit.ly", "tinyurl.com", "t.co", "goo.gl", "ow.ly", "buff.ly", "rb.gy",
    "rebrand.ly", "lnkd.in", "is.gd", "cutt.ly", "s.id", "shorturl.at",
}


# ---------------- Small utilities ----------------
def load_latest_downloader_output(env_mode: str, output_root: str) -> Tuple[Path, List[Dict[str, Any]]]:
    # Prefer pointer file first
    p = resolve_latest_pointer(env_mode, output_root, "latest_downloaded.json")
    if p is None:
        p = find_latest_matching(env_mode, output_root, "final_output_*.json")
    if p is None:
        raise RuntimeError(f"No downloader output found under {output_root}/{env_mode.upper()}/")

    payload = load_stage_payload(p)

    # Support both shapes: full payload dict OR raw docs list
    if isinstance(payload, dict):
        docs = payload.get("docs") or []
    elif isinstance(payload, list):
        docs = payload
    else:
        raise RuntimeError(f"Unexpected payload shape in {p}: {type(payload)}")

    if not isinstance(docs, list):
        raise RuntimeError(f"Invalid docs in {p}")

    return p, docs

def load_latest_enricher_output(env_mode: str, output_root: str) -> Tuple[Optional[Path], List[Dict[str, Any]]]:
    # Prefer pointer file first
    p = resolve_latest_pointer(env_mode, output_root, "latest_enriched.json")
    if p is None:
        p = find_latest_matching(env_mode, output_root, "final_enriched_*.json")
    if p is None:
        return None, []

    payload = load_stage_payload(p)
    if isinstance(payload, dict):
        docs = payload.get("docs") or []
    elif isinstance(payload, list):
        docs = payload
    else:
        return p, []

    if not isinstance(docs, list):
        return p, []

    return p, docs

def load_stage_payload(path: Path) -> Dict[str, Any]:
    with path.open("r", encoding="utf-8") as fh:
        return json.load(fh)

def safe_host(url: Any) -> str:
    try:
        from urllib.parse import urlparse
        if isinstance(url, (bytes, bytearray)):
            url = url.decode("utf-8", "ignore")
        u = urlparse(str(url))
        return f"{u.scheme}://{u.netloc}"
    except Exception:
        return "unknown://"

def _host(url: str) -> str:
    try:
        return (urlparse(url).netloc or "").lower()
    except Exception:
        return ""

def resolve_final_url_cached(url: Any, *, timeout: int = 8) -> Optional[str]:
    if not isinstance(url, str):
        return None
    u = url.strip()
    if not u:
        return None

    with _RESOLVE_LOCK:
        cached = _RESOLVE_CACHE.get(u)
    if cached is not None:
        return cached or None  # we store "" for failures

    final = resolve_final_url_http_only(u, timeout=timeout)

    with _RESOLVE_LOCK:
        _RESOLVE_CACHE[u] = final or ""

    return final

def is_youtube_host(host: str) -> bool:
    # Keep it strict: only YouTube domains
    # (youtu.be is a redirector too)
    h = (host or "").lower()
    if h.startswith("www."):
        h = h[4:]
    return h in {"youtube.com", "youtu.be"} or h.endswith(".youtube.com")

def resolve_final_url_http_only(url: Any, *, timeout: int = 8) -> Optional[str]:
    """
    Resolve HTTP redirects (not JS redirects). Returns final URL or None.
    Uses a shared public session to benefit from connection pooling.
    """
    if not isinstance(url, str):
        return None
    u = url.strip()
    if not u:
        return None

    sess = get_public_session(timeout=timeout)

    try:
        r = sess.head(u, allow_redirects=True)
        # Some servers don't support HEAD well; fall back to GET.
        if r.status_code >= 400:
            r = sess.get(u, allow_redirects=True)
        return r.url
    except Exception:
        return None

def resolves_to_youtube(url: Any) -> bool:
    final = resolve_final_url_cached(url, timeout=int(os.getenv("URL_RESOLVE_TIMEOUT", "8")))
    return bool(final and is_youtube_host(_host(final)))

def insert_after_key(d: Dict[str, Any], *, after_key: str, insert_key: str, insert_value: Any) -> None:
    """
    Rebuild dict to ensure insert_key appears immediately after after_key.
    No-op if after_key not present.
    """
    if after_key not in d:
        d[insert_key] = insert_value
        return

    new_d: Dict[str, Any] = {}
    for k, v in d.items():
        new_d[k] = v
        if k == after_key:
            new_d[insert_key] = insert_value

    d.clear()
    d.update(new_d)

def env_bool(name: str, default: bool = True) -> bool:
    v = os.getenv(name)
    if v is None:
        return default
    return v.strip().lower() in {"1", "true", "yes", "y", "on"}


def _looks_like_binary_payload_text(text: Any) -> bool:
    if not isinstance(text, str):
        return False
    if not text:
        return False
    sample = text[:4096]
    markers = (
        "[Content_Types].xml",
        "ppt/",
        "word/",
        "xl/",
        "docProps/",
        "_rels/",
        "PK\x03\x04",
    )
    if sample.startswith("PK\x03\x04") and any(marker in sample for marker in markers[:-1]):
        return True
    if any(ch in sample for ch in ("\x00", "\ufffd")):
        return True
    nonprintable = sum(1 for ch in sample if ord(ch) < 32 and ch not in "\n\r\t")
    if sample and (nonprintable / max(1, len(sample))) > 0.05:
        return True
    return False


def _looks_like_zip_bytes(raw: bytes) -> bool:
    return isinstance(raw, (bytes, bytearray)) and raw[:4] == b"PK\x03\x04"


def _office_kind_from_response(*, url: str, content_type: str, raw: bytes) -> Optional[str]:
    url_l = (url or "").lower()
    ct = (content_type or "").lower()
    if url_l.endswith(".pptx") or "presentationml.presentation" in ct:
        return "pptx"
    if url_l.endswith(".docx") or "wordprocessingml.document" in ct:
        return "docx"
    if url_l.endswith(".xlsx") or "spreadsheetml.sheet" in ct:
        return "xlsx"
    if _looks_like_zip_bytes(raw):
        try:
            with ZipFile(BytesIO(raw)) as zf:
                names = set(zf.namelist())
        except Exception:
            return None
        if any(name.startswith("ppt/") for name in names):
            return "pptx"
        if any(name.startswith("word/") for name in names):
            return "docx"
        if any(name.startswith("xl/") for name in names):
            return "xlsx"
    return None


def _extract_pptx_text(raw: bytes) -> Optional[str]:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(BytesIO(raw))
    except Exception:
        return None
    parts: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: List[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str):
                text = text.strip()
                if text:
                    slide_parts.append(text)
        if slide_parts:
            parts.append(f"Slide {idx}\n" + "\n".join(slide_parts))
    text = "\n\n".join(parts).strip()
    return text or None


def _extract_docx_text(raw: bytes) -> Optional[str]:
    try:
        with ZipFile(BytesIO(raw)) as zf:
            xml = zf.read("word/document.xml")
    except Exception:
        return None
    try:
        root = ET.fromstring(xml)
    except Exception:
        return None
    texts = [node.text.strip() for node in root.findall(".//w:t", _OOXML_NS) if isinstance(node.text, str) and node.text.strip()]
    text = "\n".join(texts).strip()
    return text or None


def _extract_xlsx_text(raw: bytes) -> Optional[str]:
    try:
        from openpyxl import load_workbook  # type: ignore
        wb = load_workbook(filename=BytesIO(raw), read_only=True, data_only=True)
    except Exception:
        return None
    parts: List[str] = []
    for ws in wb.worksheets:
        rows: List[str] = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if vals:
                rows.append(" | ".join(vals))
        if rows:
            parts.append(f"Sheet {ws.title}\n" + "\n".join(rows))
    text = "\n\n".join(parts).strip()
    return text or None


def extract_office_text_from_bytes(raw: bytes, *, kind: str) -> Optional[str]:
    if kind == "pptx":
        return _extract_pptx_text(raw)
    if kind == "docx":
        return _extract_docx_text(raw)
    if kind == "xlsx":
        return _extract_xlsx_text(raw)
    return None


def extract_office_text_from_url(url: str, *, timeout: int = 60, min_chars: int = 100) -> Tuple[Optional[str], str]:
    sess = get_public_session(timeout=timeout)
    try:
        resp = sess.get(url, allow_redirects=True, timeout=timeout)
        resp.raise_for_status()
    except Exception as e:
        logger.warning("[OfficeUrlExtractFail] host=%s url=%s reason=download err=%r", safe_host(url), url, e)
        return None, "office_url:download_failed"
    raw = resp.content or b""
    if not raw:
        return None, "office_url:empty"
    final_url = getattr(resp, "url", url) or url
    kind = _office_kind_from_response(
        url=final_url,
        content_type=resp.headers.get("content-type", ""),
        raw=raw,
    )
    if not kind:
        return None, "office_url:not_office"
    text = extract_office_text_from_bytes(raw, kind=kind)
    if not isinstance(text, str):
        logger.warning("[OfficeUrlExtractFail] host=%s url=%s reason=parse_failed kind=%s", safe_host(url), final_url, kind)
        return None, f"office_url:{kind}:parse_failed"
    text = text.strip()
    if len(text) < min_chars:
        return None, f"office_url:{kind}:too_short"
    logger.info("[OfficeUrlExtractOK] host=%s url=%s kind=%s chars=%s", safe_host(url), final_url, kind, len(text))
    return text, f"office_url:{kind}"

# ------------------------ HTTP session helpers ------------------------
def get_public_session(timeout: int = 30) -> requests.Session:
    """
    Thread-local session for calling external services (PageSense, sniffing, etc).
    Does NOT include backend basic auth.
    """
    key = f"public|{timeout}"
    sess = getattr(_tls, "session", None)
    sess_key = getattr(_tls, "session_key", None)

    if sess is None or sess_key != key:
        try:
            if sess is not None:
                sess.close()
        except Exception:
            pass

        s = requests.Session()
        s.headers.update({"accept": "application/json"})

        # Default timeout wrapper (like downloader)
        original = s.request

        def _with_timeout(method, url, **kw):
            kw.setdefault("timeout", timeout)
            return original(method, url, **kw)

        s.request = _with_timeout  # type: ignore[attr-defined]

        _tls.session = s
        _tls.session_key = key
        sess = s

    return sess


# ------------------------ pagesense runner ------------------------
def fetch_url_text_with_extractor(
    url: str,
    timeout: Optional[int] = None,
    retries: int = None,
    backoff: float = None,
    min_chars: int = None,
    session: Optional[requests.Session] = None,
) -> Optional[str]:
    """
    Calls text-extractor:
      GET {EXTRACTOR_BASE}/api/extract?url=<encoded>
    Returns raw extracted text (str) or None on failure.
    """
    base = os.getenv("URL_CONTENT_EXTRACTOR_BASE", "https://pagesense.nexavion.com").rstrip("/")
    endpoint = f"{base}/api/extract"
    timeout = int(os.getenv("EXTRACTOR_TIMEOUT", str(timeout or 720)))
    retries = int(os.getenv("EXTRACTOR_RETRIES", str(retries or 3)))
    backoff = float(os.getenv("EXTRACTOR_BACKOFF", str(backoff or 1.6)))
    min_chars = int(os.getenv("EXTRACTOR_MIN_CHARS", str(min_chars or 100)))

    sess = session or requests.Session()

    last_err = None
    for attempt in range(retries + 1):
        try:
            # small pre-call jitter to avoid thundering herd
            time.sleep(random.uniform(0.15, 0.45))

            r = sess.get(endpoint, params={"url": url}, timeout=timeout)
            if r.status_code == 429 or 500 <= r.status_code < 600:
                raise requests.HTTPError(f"HTTP {r.status_code}", response=r)
            if r.status_code != 200:
                return None

            body = r.json() if r.headers.get("content-type", "").startswith("application/json") else {}
            if not body or not body.get("ok"):
                return None

            text = (body.get("text") or "").strip()
            if len(text) < min_chars:
                return None
            if _looks_like_binary_payload_text(text):
                logger.warning(
                    "[ExtractorBinaryPayload] host=%s url=%s len=%s",
                    safe_host(url),
                    url,
                    len(text),
                )
                return None
            return text

        except (requests.Timeout, requests.ConnectionError, requests.HTTPError) as e:
            last_err = e
            if attempt < retries:
                # exponential backoff + jitter
                delay = (backoff ** attempt) + random.uniform(0.2, 0.6)
                time.sleep(delay)
                continue
            return None
        except Exception as e:
            last_err = e
            return None

# def pagesense_one(llid: str, url: str) -> Tuple[str, str, Optional[str], str]:
def pagesense_one(*, llid: str, url: str, http_timeout: int) -> Tuple[str, str, Optional[str], str]:
    time.sleep(random.uniform(0.05, 0.25))
    try:
        sess = get_public_session(timeout=int(os.getenv("EXTRACTOR_HTTP_TIMEOUT", "35")))
        text = fetch_url_text_with_extractor(
            url,
            timeout=int(os.getenv("EXTRACTOR_TIMEOUT", "150")),
            retries=int(os.getenv("EXTRACTOR_RETRIES", "3")),
            backoff=float(os.getenv("EXTRACTOR_BACKOFF", "1.6")),
            min_chars=int(os.getenv("EXTRACTOR_MIN_CHARS", "100")),
            session=sess,
        )
        return llid, url, (text.strip() if isinstance(text, str) else None), "pagesense"
    except Exception as e:
        logger.error("[PageSenseFail] id=%s host=%s err=%s", llid, safe_host(url), e)
        return llid, url, None, "pagesense:failed"


# ------------------------ transcribe with runpod custom runner ------------------------
def transcribe_with_custom_api(*, media_url: str) -> Optional[str]:
    """
    Calls the RunPod /transcribe endpoint with a direct media URL (S3 hosted audio/video).
    Returns transcript text or None.
    """
    user = (os.getenv("CUSTOM_TRANSCRIBE_USER") or "").strip()
    pwd = (os.getenv("CUSTOM_TRANSCRIBE_PASSWORD") or "").strip()

    auth = None
    if user and pwd:
        auth = HTTPBasicAuth(user, pwd)

    retries = int(os.getenv("CUSTOM_TRANSCRIBE_RETRIES", "2"))
    backoff = float(os.getenv("CUSTOM_TRANSCRIBE_BACKOFF", "1.5"))
    endpoint = (os.getenv("CUSTOM_TRANSCRIBE_ENDPOINT") or "").strip()

    if not endpoint:
        logger.error("[CustomTranscribe] Missing CUSTOM_TRANSCRIBE_ENDPOINT env var")
        return None

    whisper_model = (os.getenv("CUSTOM_TRANSCRIBE_MODEL") or "large-v1").strip() or "large-v1"
    timeout_s = float(os.getenv("CUSTOM_TRANSCRIBE_HTTP_TIMEOUT", "180"))  # transcription can take a while

    for attempt in range(retries + 1):
        try:
            sess = get_public_session(timeout=int(timeout_s))
            r = sess.post(
                endpoint.rstrip("/") + "/transcribe",
                json={"url": media_url, "whisper_model": whisper_model},
                headers={"accept": "application/json", "content-type": "application/json"},
                auth=auth,
            )

            if r.status_code >= 400:
                logger.error(
                    "[CustomTranscribeFail] media_host=%s status=%s body=%s",
                    media_url,
                    r.status_code,
                    (r.text or "")[:500],  # cap to avoid huge logs
                )
                return None

            r.raise_for_status()
            data = r.json()

            text = (((data or {}).get("whisper") or {}).get("text") or "")
            text = text.strip() if isinstance(text, str) else ""
            return text or None
        except Exception as e:
            if attempt >= retries:
                logger.error("[CustomTranscribeFail] host=%s err=%s", safe_host(media_url), e)
                return None
            time.sleep(backoff ** attempt)

def custom_transcribe_one(llid: str, url: str) -> Tuple[str, str, Optional[str], str]:
    """
    custom_transcribe: send direct hosted media URL to the custom endpoint.
    """
    time.sleep(random.uniform(0.05, 0.25))

    text = transcribe_with_custom_api(media_url=url)
    if text:
        return llid, url, text, "custom_transcribe"
    return llid, url, None, "custom_transcribe:failed"


def probe_target_url(url: Any, *, timeout: int = 8) -> Tuple[bool, str]:
    """
    Lightweight reachability probe for expensive enrichment routes.

    Returns:
      (True, "ok")
      (False, "<reason>")
    """
    if not isinstance(url, str) or not url.strip():
        return False, "empty"

    u = url.strip()
    ok, reason, resolved = validate_url_safety(u, follow_redirects=True, timeout=timeout, check_reputation=True)
    if not ok:
        return False, reason
    target = resolved or u
    sess = get_public_session(timeout=timeout)

    try:
        r = sess.head(target, allow_redirects=True)
        status = r.status_code
        if status == 405:
            r = sess.get(target, allow_redirects=True, stream=True)
            status = r.status_code
        try:
            r.close()
        except Exception:
            pass

        if 200 <= status < 400:
            return True, "ok"
        if status == 404:
            return False, "http_404"
        if status == 403:
            return False, "http_403"
        if 400 <= status < 500:
            return False, f"http_{status}"
        if 500 <= status < 600:
            return False, f"http_{status}"
        return False, f"http_{status}"
    except requests.Timeout:
        return False, "timeout"
    except requests.ConnectionError:
        return False, "connection_error"
    except Exception:
        return False, "probe_error"


def probe_remote_content_type(url: Any, *, timeout: int = 8) -> Optional[str]:
    if not isinstance(url, str) or not url.strip():
        return None

    u = url.strip()
    ok, _, resolved = validate_url_safety(u, follow_redirects=True, timeout=timeout, check_reputation=False)
    if not ok:
        return None
    target = resolved or u
    sess = get_public_session(timeout=timeout)

    try:
        r = sess.head(target, allow_redirects=True)
        status = r.status_code
        if status == 405 or status >= 400:
            r = sess.get(target, allow_redirects=True, stream=True)
        content_type = (r.headers.get("content-type") or "").split(";", 1)[0].strip().lower() or None
        try:
            r.close()
        except Exception:
            pass
        return content_type
    except Exception:
        return None


def probe_media_duration_seconds(url: Any, *, timeout: int = 20) -> Optional[float]:
    if not isinstance(url, str) or not url.strip():
        return None

    cmd = [
        "/usr/bin/ffprobe",
        "-v",
        "error",
        "-show_entries",
        "format=duration",
        "-of",
        "default=noprint_wrappers=1:nokey=1",
        url.strip(),
    ]

    try:
        proc = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=max(timeout, 5),
        )
    except Exception as e:
        logger.info("[MediaDurationProbeFail] host=%s err=%r", safe_host(url), e)
        return None

    if proc.returncode != 0:
        logger.info(
            "[MediaDurationProbeFail] host=%s rc=%s stderr=%s",
            safe_host(url),
            proc.returncode,
            (proc.stderr or "")[:300],
        )
        return None

    try:
        value = float((proc.stdout or "").strip())
    except Exception:
        return None

    return value if value > 0 else None


def transcribe_with_custom_api_detailed(*, media_url: str) -> Tuple[Optional[str], str]:
    """
    Calls the RunPod /transcribe endpoint and returns (text, tag).
    """
    user = (os.getenv("CUSTOM_TRANSCRIBE_USER") or "").strip()
    pwd = (os.getenv("CUSTOM_TRANSCRIBE_PASSWORD") or "").strip()

    auth = None
    if user and pwd:
        auth = HTTPBasicAuth(user, pwd)

    retries = int(os.getenv("CUSTOM_TRANSCRIBE_RETRIES", "2"))
    backoff = float(os.getenv("CUSTOM_TRANSCRIBE_BACKOFF", "1.5"))
    endpoint = (os.getenv("CUSTOM_TRANSCRIBE_ENDPOINT") or "").strip()

    if not endpoint:
        logger.error("[CustomTranscribe] Missing CUSTOM_TRANSCRIBE_ENDPOINT env var")
        return None, "custom_transcribe:missing_endpoint"

    max_duration_sec = float(os.getenv("CUSTOM_TRANSCRIBE_MAX_DURATION_SEC", "3000") or "3000")
    duration_probe_timeout = int(os.getenv("CUSTOM_TRANSCRIBE_DURATION_PROBE_TIMEOUT", "20"))
    if max_duration_sec > 0:
        content_type = probe_remote_content_type(media_url, timeout=min(duration_probe_timeout, 10))
        if isinstance(content_type, str) and content_type.startswith("video/"):
            duration_sec = probe_media_duration_seconds(media_url, timeout=duration_probe_timeout)
            if duration_sec is not None and duration_sec > max_duration_sec:
                logger.warning(
                    "[CustomTranscribeSkip] media_host=%s content_type=%s duration_sec=%.2f max_sec=%.2f",
                    media_url,
                    content_type,
                    duration_sec,
                    max_duration_sec,
                )
                return None, "custom_transcribe:too_long"

    whisper_model = (os.getenv("CUSTOM_TRANSCRIBE_MODEL") or "large-v1").strip() or "large-v1"
    timeout_s = float(os.getenv("CUSTOM_TRANSCRIBE_HTTP_TIMEOUT", "180"))

    for attempt in range(retries + 1):
        try:
            sess = get_public_session(timeout=int(timeout_s))
            r = sess.post(
                endpoint.rstrip("/") + "/transcribe",
                json={"url": media_url, "whisper_model": whisper_model},
                headers={"accept": "application/json", "content-type": "application/json"},
                auth=auth,
            )

            if r.status_code >= 400:
                logger.error(
                    "[CustomTranscribeFail] media_host=%s status=%s body=%s",
                    media_url,
                    r.status_code,
                    (r.text or "")[:500],
                )
                if r.status_code == 524:
                    return None, "custom_transcribe:proxy_timeout"
                if r.status_code == 404:
                    return None, "custom_transcribe:http_404"
                if r.status_code == 429:
                    return None, "custom_transcribe:http_429"
                if 500 <= r.status_code < 600:
                    return None, f"custom_transcribe:http_{r.status_code}"
                return None, f"custom_transcribe:http_{r.status_code}"

            r.raise_for_status()
            data = r.json()
            text = (((data or {}).get("whisper") or {}).get("text") or "")
            text = text.strip() if isinstance(text, str) else ""
            return (text or None), ("custom_transcribe" if text else "custom_transcribe:empty")
        except requests.Timeout:
            if attempt >= retries:
                logger.error("[CustomTranscribeFail] host=%s err=timeout", safe_host(media_url))
                return None, "custom_transcribe:timeout"
            time.sleep(backoff ** attempt)
        except Exception as e:
            if attempt >= retries:
                logger.error("[CustomTranscribeFail] host=%s err=%s", safe_host(media_url), e)
                return None, "custom_transcribe:failed"
            time.sleep(backoff ** attempt)

    return None, "custom_transcribe:failed"

def _normalise_netloc(netloc: str) -> str:
    n = (netloc or "").strip().lower()
    # Drop common prefix for checks
    if n.startswith("www."):
        n = n[4:]
    return n


def _url_safety_enabled() -> bool:
    return env_bool("URL_SAFETY_ENABLE", True)


def _safe_browsing_enabled() -> bool:
    return env_bool("URL_SAFETY_ENABLE_SAFE_BROWSING", False) and bool((os.getenv("GOOGLE_SAFE_BROWSING_API_KEY") or "").strip())


def _is_shortener_host(host: str) -> bool:
    return _normalise_netloc(host) in _SHORTENER_HOSTS


def _host_is_obviously_internal(host: str) -> bool:
    h = _normalise_netloc(host)
    if not h:
        return True
    if h in {"localhost", "localhost.localdomain"}:
        return True
    if h.endswith(".local") or h.endswith(".localdomain") or h.endswith(".internal") or h.endswith(".home") or h.endswith(".lan"):
        return True
    try:
        ip = ipaddress.ip_address(h)
    except Exception:
        return False
    return (
        ip.is_private
        or ip.is_loopback
        or ip.is_link_local
        or ip.is_multicast
        or ip.is_reserved
        or ip.is_unspecified
    )


def _host_resolves_internal(host: str) -> bool:
    h = _normalise_netloc(host)
    if _host_is_obviously_internal(h):
        return True
    try:
        infos = socket.getaddrinfo(h, None, type=socket.SOCK_STREAM)
    except Exception as e:
        logger.info("[UrlSafetyResolveSkip] host=%s err=%r", h, e)
        return False
    seen = False
    for info in infos:
        sockaddr = info[4]
        if not sockaddr:
            continue
        addr = sockaddr[0]
        try:
            ip = ipaddress.ip_address(addr)
        except Exception:
            continue
        seen = True
        if (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_reserved
            or ip.is_unspecified
        ):
            return True
    return False if not seen else False


def _path_has_dangerous_extension(path: str) -> bool:
    p = (path or "").strip().lower()
    if not p:
        return False
    for ext in _LEGACY_WEBPAGE_EXTENSIONS:
        if p.endswith(ext):
            return False
    return any(p.endswith(ext) for ext in _DANGEROUS_EXTENSIONS)


def _safe_browsing_matches(url: str) -> bool:
    api_key = (os.getenv("GOOGLE_SAFE_BROWSING_API_KEY") or "").strip()
    if not api_key:
        return False
    with _SAFE_BROWSING_LOCK:
        cached = _SAFE_BROWSING_CACHE.get(url)
    if cached is not None:
        return cached
    endpoint = f"https://safebrowsing.googleapis.com/v4/threatMatches:find?key={api_key}"
    payload = {
        "client": {"clientId": "data-prep-opensearch", "clientVersion": "1.0"},
        "threatInfo": {
            "threatTypes": [
                "MALWARE",
                "SOCIAL_ENGINEERING",
                "UNWANTED_SOFTWARE",
                "POTENTIALLY_HARMFUL_APPLICATION",
            ],
            "platformTypes": ["ANY_PLATFORM"],
            "threatEntryTypes": ["URL"],
            "threatEntries": [{"url": url}],
        },
    }
    try:
        sess = get_public_session(timeout=int(os.getenv("URL_SAFETY_SAFE_BROWSING_TIMEOUT", "10")))
        r = sess.post(endpoint, json=payload, headers={"accept": "application/json", "content-type": "application/json"})
        r.raise_for_status()
        data = r.json() if r.headers.get("content-type", "").startswith("application/json") else {}
        flagged = bool((data or {}).get("matches"))
    except Exception as e:
        logger.warning("[UrlSafetySafeBrowsingFail] url=%s err=%r", url, e)
        flagged = False
    with _SAFE_BROWSING_LOCK:
        _SAFE_BROWSING_CACHE[url] = flagged
    return flagged


def _check_single_url_safety(url: str, *, check_reputation: bool) -> Tuple[bool, str]:
    if not isinstance(url, str) or not url.strip():
        return False, "empty"
    u = url.strip()
    try:
        p = urlparse(u)
    except Exception:
        return False, "parse_error"
    scheme = (p.scheme or "").lower()
    if scheme != "https":
        return False, "https_required"
    if p.username or p.password:
        return False, "embedded_credentials"
    host = _normalise_netloc(p.hostname or p.netloc)
    if not host:
        return False, "missing_host"
    if _host_is_obviously_internal(host):
        return False, "internal_host"
    if "." not in host:
        return False, "host_missing_tld"
    tld = host.rsplit(".", 1)[-1]
    if len(tld) < 2:
        return False, "host_bad_tld"
    if _path_has_dangerous_extension(p.path or ""):
        return False, "dangerous_payload_extension"
    if check_reputation and _safe_browsing_enabled() and _safe_browsing_matches(u):
        return False, "safe_browsing_flagged"
    return True, "ok"


def validate_url_safety(url: Any, *, follow_redirects: bool, timeout: int, check_reputation: bool) -> Tuple[bool, str, Optional[str]]:
    if isinstance(url, (bytes, bytearray)):
        url = url.decode("utf-8", "ignore")
    if not isinstance(url, str):
        return False, "not_a_string", None
    original = url.strip()
    if not original:
        return False, "empty", None
    if not _url_safety_enabled():
        return True, "ok", original
    ok, reason = _check_single_url_safety(original, check_reputation=check_reputation)
    if not ok:
        return False, reason, None
    try:
        p = urlparse(original)
        host = _normalise_netloc(p.hostname or p.netloc)
        if _host_resolves_internal(host):
            return False, "internal_resolution", None
    except Exception:
        return False, "parse_error", None
    if not follow_redirects:
        return True, "ok", original
    final = resolve_final_url_http_only(original, timeout=timeout)
    if not final:
        return False, "unreachable_or_unresolved", None
    ok, reason = _check_single_url_safety(final, check_reputation=check_reputation)
    if not ok:
        return False, f"final_{reason}", None
    try:
        p = urlparse(final)
        host = _normalise_netloc(p.hostname or p.netloc)
        if _host_resolves_internal(host):
            return False, "final_internal_resolution", None
    except Exception:
        return False, "final_parse_error", None
    if _is_shortener_host(urlparse(original).hostname or "") and check_reputation and _safe_browsing_enabled():
        if _safe_browsing_matches(original):
            return False, "shortener_safe_browsing_flagged", None
        if _safe_browsing_matches(final):
            return False, "final_safe_browsing_flagged", None
    return True, "ok", final

def classify_url_for_enrichment(url: Any) -> Tuple[bool, str]:
    """
    Returns (ok, reason).
    ok=True  -> safe to send to pagesense/transcribe
    ok=False -> skip with 'reason'
    """
    if isinstance(url, (bytes, bytearray)):
        url = url.decode("utf-8", "ignore")

    if not isinstance(url, str):
        return False, "not_a_string"

    u = url.strip()
    if not u:
        return False, "empty"

    try:
        p = urlparse(u)
    except Exception:
        return False, "parse_error"

    scheme = (p.scheme or "").lower()
    if scheme != "https":
        return False, "https_required"

    if p.username or p.password:
        return False, "embedded_credentials"

    netloc = _normalise_netloc(p.netloc)
    if not netloc:
        return False, "missing_host"
    if _host_is_obviously_internal(p.hostname or netloc):
        return False, "internal_host"

    path = p.path or ""
    path_l = path.lower()

    if netloc == "github.com":
        if path_l.startswith("/orgs/") and "/projects/" in path_l:
            return False, "github_project_board"
        if re.match(r"^/[^/]+/[^/]+/projects(?:/|$)", path_l):
            return False, "github_project_board"
        if path_l in {"/notifications", "/pulls", "/issues", "/settings"}:
            return False, "github_dashboard_page"

    # Reject obviously broken hosts like "www.google" (no TLD)
    # Heuristic: must contain at least one dot and end segment length >= 2
    if "." not in netloc:
        return False, "host_missing_tld"
    tld = netloc.rsplit(".", 1)[-1]
    if len(tld) < 2:
        return False, "host_bad_tld"

    if _path_has_dangerous_extension(path):
        return False, "dangerous_payload_extension"

    # Homepage / bare domain: no real path (or "/") AND no query params
    has_query = bool(p.query and p.query.strip())
    if (path == "" or path == "/") and not has_query:
        return False, "homepage_or_bare_domain"

    ok, reason, _ = validate_url_safety(
        u,
        follow_redirects=True,
        timeout=int(os.getenv("URL_RESOLVE_TIMEOUT", "8")),
        check_reputation=True,
    )
    if not ok:
        return False, reason

    return True, "ok"



def target_url_for_enrichment(doc: Dict[str, Any]) -> Optional[str]:
    """
    Decide what URL we enrich from.

    - pagesense: always doc["@id"]
    - transcribe:
        - if hosted media: use ko_file_id (S3 hosted URL)
        - else: use @id (youtube/vimeo/etc.)
    """
    via = doc.get("enrich_via")
    at_id = doc.get("@id") if isinstance(doc.get("@id"), str) else None
    ko_file_id = doc.get("ko_file_id") if isinstance(doc.get("ko_file_id"), str) else None
    ko_is_hosted = bool(doc.get("ko_is_hosted"))

    if via == "pagesense":
        if ko_is_hosted and ko_file_id and ko_file_id.strip():
            return ko_file_id.strip()
        return at_id.strip() if at_id and at_id.strip() else None

    mimetype = (doc.get("ko_object_mimetype") or "").strip().lower()

    if via == "api_transcribe":
        # hosted media (video/audio/image) -> use hosted file url
        if (not ko_is_hosted) and any(mimetype.startswith(x) for x in ("video/", "audio/")):
            return at_id.strip() if at_id and at_id.strip() else None

        # Non-hosted media platforms -> use @id (YouTube etc.)
        return at_id.strip() if at_id and at_id.strip() else None

    if via == "custom_transcribe":
        # hosted media (video/audio/image) -> use hosted file url
        if ko_is_hosted and any(mimetype.startswith(x) for x in ("video/", "audio/")):
            return ko_file_id.strip() if ko_file_id and ko_file_id.strip() else None

        # non-hosted media platforms -> use @id
        return at_id.strip() if at_id and at_id.strip() else None

    return None

def is_media_mimetype(m: Any) -> bool:
    """
    Treat image/*, audio/*, video/* as media requiring transcription.
    """
    if not isinstance(m, str):
        return False
    mm = m.strip().lower()
    # mm.startswith("image/") or
    return mm.startswith("audio/") or mm.startswith("video/")


def is_visual_document_mimetype(m: Any) -> bool:
    if not isinstance(m, str):
        return False
    mm = m.strip().lower()
    return mm.startswith("image/") or mm == "application/pdf"


def as_bool(v: Any) -> bool:
    """Normalise common truthy/falsey representations."""
    if isinstance(v, bool):
        return v
    if v is None:
        return False
    if isinstance(v, (int, float)):
        return v != 0
    if isinstance(v, str):
        s = v.strip().lower()
        if s in {"1", "true", "yes", "y", "on"}:
            return True
        if s in {"0", "false", "no", "n", "off", ""}:
            return False
    # last resort: Python truthiness
    return bool(v)

def set_enrich_via(d: Dict[str, Any]) -> None:
    ko_is_hosted = as_bool(d.get("ko_is_hosted"))
    mimetype = (d.get("ko_object_mimetype") or "").strip() if isinstance(d.get("ko_object_mimetype"), str) else ""
    ko_file_id = (d.get("ko_file_id") or "").strip() if isinstance(d.get("ko_file_id"), str) else ""
    at_id = d.get("@id")

    # 1) Hosted + media mimetype => custom_transcribe
    if ko_is_hosted and is_media_mimetype(mimetype):
        d["enrich_via"] = "custom_transcribe"
        return

    # Hosted image/PDF resources should still be enriched through the visual path.
    if ko_is_hosted and is_visual_document_mimetype(mimetype):
        d["enrich_via"] = "pagesense"
        return

    # Everything below is only for non-hosted URLs
    if not ko_is_hosted:
        # inside set_enrich_via(d)
        ko_is_hosted = as_bool(d.get("ko_is_hosted"))
        at_id = d.get("@id")

        # Store resolved_url (non-hosted only), and keep it right after @id
        resolved_url: Optional[str] = None
        if isinstance(at_id, str) and at_id.strip():
            resolved_url = resolve_final_url_cached(at_id, timeout=int(os.getenv("URL_RESOLVE_TIMEOUT", "8")))
            if isinstance(resolved_url, str) and resolved_url.strip():
                d["resolved_url"] = resolved_url.strip()
                insert_after_key(d, after_key="@id", insert_key="resolved_url", insert_value=d["resolved_url"])
            else:
                d.pop("resolved_url", None)  # keep output clean

        # Optional: only do this expensive check when we otherwise would choose pagesense.
        # Focus only on YouTube, as requested.
        resolve_enabled = env_bool("ENABLE_URL_RESOLVE_YOUTUBE", default=True)

        # 2) URL-only case (non-hosted + no mimetype + no file id)
        if mimetype == "" and ko_file_id == "":
            if is_deapi_platform_url(at_id):
                d["enrich_via"] = "api_transcribe"
                return

            # Otherwise, try to detect if it *redirects* to YouTube.
            if resolve_enabled:
                final_for_check = resolved_url or None
                if final_for_check and is_deapi_platform_url(final_for_check):
                    d["enrich_via"] = "api_transcribe"
                    return

            d["enrich_via"] = "pagesense"
            return

        # 3) Non-hosted (general fallback)
        if is_deapi_platform_url(at_id):
            d["enrich_via"] = "api_transcribe"
            return

        d["enrich_via"] = "pagesense"
        return

    # Hosted but non-media: no external enrichment route needed
    d.pop("enrich_via", None)

def has_enrich_via(doc: Dict[str, Any]) -> bool:
    return doc.get("enrich_via") in {"pagesense", "api_transcribe", "custom_transcribe"}

def is_placeholder_content(val: Any) -> bool:
    if not isinstance(val, str):
        return True
    if _looks_like_binary_payload_text(val):
        return True
    return val.strip().lower() in ("", "no content present")

def already_enriched(prev_doc: Dict[str, Any]) -> bool:
    return prev_doc.get("enriched") == 1 and not is_placeholder_content(prev_doc.get("ko_content_flat"))

def should_skip(current: Dict[str, Any], prev_doc: Optional[Dict[str, Any]]) -> bool:
    """
    Skip re-enrichment iff:
      - previous doc exists
      - previous doc has actual enrichment (enriched=1 and ko_content_flat not placeholder)
      - and enricher inputs fingerprint is unchanged
    """
    if not prev_doc:
        return False

    if not already_enriched(prev_doc):
        return False

    prev_fp = prev_doc.get("_enrich_inputs_fp")
    cur_fp = current.get("_enrich_inputs_fp")

    # If either side lacks fp, fall back to conservative behaviour (do NOT skip).
    if not isinstance(prev_fp, str) or not isinstance(cur_fp, str):
        return False

    return prev_fp == cur_fp

def is_deapi_platform_url(u: str) -> bool:
    """
    deAPI /vid2txt only accepts a few platforms (per error message).
    Keep this strict to avoid 422 spam.
    """
    if not isinstance(u, str):
        return False
    v = u.lower()
    return any(host in v for host in (
        "youtube.com/", "youtu.be/",
        "twitter.com/", "x.com/",
        "twitch.tv/",
        "kick.com/",
    ))

def compute_enrich_inputs_fp(doc: Dict[str, Any]) -> str:
    """
    Compute enricher-input fingerprint according to the 3 regimes we agreed.

    A) URL-only (ko_is_hosted False) -> @id only
    B) Hosted doc-like (hosted, non media, no enrich_via) -> @id + ko_file_id + ko_content_flat
       (even if enricher does nothing, this FP is still useful for consistency)
    C) Hosted media (video/audio/image) -> @id + ko_file_id
    """
    ko_is_hosted = as_bool(doc.get("ko_is_hosted"))
    mimetype = (doc.get("ko_object_mimetype") or "").strip().lower()
    at_id = stable_str(doc.get("@id"))
    ko_file_id = stable_str(doc.get("ko_file_id"))
    ko_content_flat = stable_str(doc.get("ko_content_flat"))

    is_media = any(mimetype.startswith(x) for x in ("video/", "audio/", "image/"))

    if not ko_is_hosted:
        obj = {"mode": "url_only", "@id": at_id}
        return sha256_obj(obj)

    if is_media:
        obj = {"mode": "hosted_media", "@id": at_id, "ko_file_id": ko_file_id}
        return sha256_obj(obj)

    # hosted, non-media
    obj = {"mode": "hosted_doc", "@id": at_id, "ko_file_id": ko_file_id, "ko_content_flat": ko_content_flat}
    return sha256_obj(obj)
